# coding:utf-8

import pptx
# import pdfkit

p = pptx.Presentation()  # 生成ppt对象

layout = p.slide_layouts[1]  # 选择布局

slide = p.slides.add_slide(layout)

p.save('test1.ppt')

# pdfkit.from_string('Hello!'.encode(), 'out.pdf')
